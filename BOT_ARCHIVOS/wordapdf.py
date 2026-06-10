import os
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import pandas as pd
import win32com.client  # Solo para archivos .doc en Windows


def convertir_doc_a_docx(ruta_archivo):
    """
    Convierte un archivo .doc a .docx utilizando Microsoft Word.
    :param ruta_archivo: Ruta completa del archivo .doc.
    :return: Ruta del archivo convertido (.docx).
    """
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False  # Mantener la ventana oculta
    try:
        doc = word.Documents.Open(ruta_archivo)
        ruta_convertida = ruta_archivo.replace('.doc', '.docx')
        doc.SaveAs(ruta_convertida, FileFormat=16)  # 16 es el formato para .docx
        doc.Close()
        print(f"Archivo convertido de .doc a .docx: {ruta_convertida}")
        return ruta_convertida
    except Exception as e:
        print(f"Error al convertir el archivo .doc: {e}")
    finally:
        word.Quit()


def convertir_docx_a_pdf(ruta_archivo, pdf_output):
    """
    Convierte un archivo .docx en un PDF.
    :param ruta_archivo: Ruta completa del archivo .docx.
    :param pdf_output: Ruta completa del archivo de salida .pdf.
    """
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False  # Ejecutar en segundo plano
    try:
        # Convertir rutas relativas a absolutas
        ruta_archivo = os.path.abspath(ruta_archivo)
        pdf_output = os.path.abspath(pdf_output)

        # Abrir y guardar el archivo
        doc = word.Documents.Open(ruta_archivo)
        doc.SaveAs(pdf_output, FileFormat=17)  # 17 es el formato para guardar como PDF
        doc.Close()
        print(f"Archivo convertido a PDF: {pdf_output}")
    except Exception as e:
        print(f"Error al convertir {ruta_archivo} a PDF: {e}")
    finally:
        word.Quit()



def convertir_txt_a_pdf(ruta_archivo, pdf_output):
    """
    Convierte un archivo de texto (.txt) en un PDF.
    :param ruta_archivo: Ruta completa del archivo .txt.
    :param pdf_output: Ruta completa del archivo de salida .pdf.
    """
    pdf = canvas.Canvas(pdf_output, pagesize=letter)
    pdf.setFont("Helvetica", 12)
    with open(ruta_archivo, 'r', encoding='utf-8') as file:
        y = 750
        for line in file:
            pdf.drawString(50, y, line.strip())
            y -= 15
            if y < 50:
                pdf.showPage()
                y = 750
    pdf.save()
    print(f"Archivo convertido a PDF: {pdf_output}")


def convertir_excel_a_pdf(ruta_archivo, pdf_output):
    """
    Convierte un archivo Excel (.xlsx) en un PDF.
    :param ruta_archivo: Ruta completa del archivo .xlsx.
    :param pdf_output: Ruta completa del archivo de salida .pdf.
    """
    data = pd.read_excel(ruta_archivo)
    pdf = canvas.Canvas(pdf_output, pagesize=letter)
    pdf.setFont("Helvetica", 10)
    y = 750
    for index, row in data.iterrows():
        line = ' '.join(map(str, row.values))
        pdf.drawString(50, y, line)
        y -= 15
        if y < 50:
            pdf.showPage()
            y = 750
    pdf.save()
    print(f"Archivo convertido a PDF: {pdf_output}")


def convertir_pptx_a_pdf(ruta_archivo, pdf_output):
    """
    Convierte un archivo PowerPoint (.pptx) en un PDF.
    :param ruta_archivo: Ruta completa del archivo .pptx.
    :param pdf_output: Ruta completa del archivo de salida .pdf.
    """
    pdf = canvas.Canvas(pdf_output, pagesize=letter)
    ppt = Presentation(ruta_archivo)
    y = 750
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.drawString(50, y, shape.text)
                y -= 15
                if y < 50:
                    pdf.showPage()
                    y = 750
    pdf.save()
    print(f"Archivo convertido a PDF: {pdf_output}")


# Programa principal para probar las conversiones
def main():
    carpeta_pruebas = "./ORTIZ.014.24/01SolicitudAudiencia"  # Ruta donde se encuentran los archivos de prueba
    archivos = os.listdir(carpeta_pruebas)

    for archivo in archivos:
        ruta_archivo = os.path.join(carpeta_pruebas, archivo)
        pdf_output = os.path.splitext(ruta_archivo)[0] + ".pdf"

        try:
            # Manejar archivos .doc convirtiéndolos a .docx
            if archivo.lower().endswith('.doc'):
                ruta_archivo = convertir_doc_a_docx(ruta_archivo)
                os.remove(os.path.join(carpeta_pruebas, archivo))  # Eliminar el archivo .doc original

            # Manejar archivos Word (.docx)
            if archivo.lower().endswith('.docx'):
                convertir_docx_a_pdf(ruta_archivo, pdf_output)

            # Manejar archivos de texto
            elif archivo.lower().endswith('.txt'):
                convertir_txt_a_pdf(ruta_archivo, pdf_output)

            # Manejar archivos Excel
            elif archivo.lower().endswith('.xlsx'):
                convertir_excel_a_pdf(ruta_archivo, pdf_output)

            # Manejar archivos PowerPoint
            elif archivo.lower().endswith('.pptx'):
                convertir_pptx_a_pdf(ruta_archivo, pdf_output)

            print(f"Archivo convertido: {pdf_output}")

        except Exception as e:
            print(f"Error al procesar el archivo {ruta_archivo}: {e}")


if __name__ == "__main__":
    main()
